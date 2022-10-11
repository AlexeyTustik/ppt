from pptx import Presentation
from pptx.util import Pt
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
import json


def add_text_box(slide, text, left, top, width=None, height=None, font_size=28):
    text_box = slide.shapes.add_textbox(
        Pt(left), Pt(top), Pt(width), Pt(height))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)


def add_table(slide, data: list, left, top, width, height):
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols, Pt(
        left), Pt(top), Pt(width), Pt(height)).table
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = str(data[i][j])
    return table


def add_picture(slide, path, left, top, width=None, height=None):
    slide.shapes.add_picture(path, Pt(left), Pt(top), Pt(width), Pt(height))


def add_bar_plot(slide, categories, series, left, top, width, height):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('Series 1', series)

    # add chart to slide --------------------
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Pt(left), Pt(
            top), Pt(width), Pt(height), chart_data
    )


def add_pie_plot(slide, categories, series, left, top, width, height):
    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series('Series 1', series)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Pt(left), Pt(top), Pt(width), Pt(height), chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    return chart


def make_test_presentation():
    # make presentation file
    prs = Presentation()
    prs.slide_width = Pt(1280)
    prs.slide_height = Pt(720)
    # 6 - blank layout
    # make slide
    slide_1 = prs.slides.add_slide(prs.slide_layouts[6])

    header_text = '''Исполнение КПЭ
    «Повышение исполняемости договоров», от 100 млн руб.*'''
    add_text_box(slide_1, header_text, 10, 10, 900, 20)
    add_picture(slide_1, 'src/logo.png', 1280-250, 10, 200, 75)

    data = [
        [
            'Общее количество КС до конца года*',
            'Общее количество КС на дату',
            'Количество КС исполненных в срок',
            'Количество КС  не исполненных в срок',
            'Фактическое значение доли КС, выполненных в срок, %'],
        [100, 100, 100, 100, 100]
    ]
    add_table(slide_1, data, 50, 100, 800, 200)
    categories = ['АЭМ', 'ЯОК', 'КРЭА', 'АСЭ', 'АРМЗ', 'ТВЭЛ', 'Итого']
    series = [66.67,	75,	75,	83.3,	84.6,	87.5,	81.82]
    add_bar_plot(slide_1, categories, series, 50, 350, 800, 250)

    categories_pie = ['Всего', 'Со сроком на 25.07',
                      'Исполнено', 'Неисполнено']
    series = [184, 88, 61, 21]
    add_pie_plot(slide_1, categories_pie, series, 900, 100, 350, 350)
    prs.save('out/new.pptx')


if __name__ == '__main__':
    make_test_presentation()
    # 123
